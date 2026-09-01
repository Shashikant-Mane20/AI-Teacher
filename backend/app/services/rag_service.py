from io import BytesIO
from typing import List

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


class RAGService:
    async def extract_text(self, file_name: str, content: bytes) -> str:
        extension = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
        if extension == "pdf":
            reader = PdfReader(BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages).strip()
        if extension == "docx":
            document = Document(BytesIO(content))
            return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()
        if extension == "pptx":
            presentation = Presentation(BytesIO(content))
            slides = []
            for slide in presentation.slides:
                slides.append("\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")))
            return "\n\n".join(slides).strip()
        if extension in {"txt", "md"}:
            return content.decode("utf-8", errors="replace").strip()
        raise ValueError("Supported file types are PDF, DOCX, PPTX, TXT, and MD")

    async def index_document(self, text: str):
        chunks = self.chunk_text(text)
        return {"chunks_count": len(chunks), "chunks": chunks}

    def chunk_text(self, text: str, chunk_size: int = 1000) -> List[str]:
        if not text:
            return []
        return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]
